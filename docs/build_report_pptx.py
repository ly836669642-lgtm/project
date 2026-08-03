"""Generates docs/Report.pptx: a 6-slide, ~7-minute talk summarizing
DOCUMENTATION.md (cover, system overview, perception, planning, control,
results). Content and figures are pulled directly from DOCUMENTATION.md /
DOCUMENTATION.tex; regenerate this script's output if the doc changes.

Run with: conda run -n claude python docs/build_report_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

DOCS = Path(__file__).parent
FIGS = DOCS / "figures"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x2E, 0x6F, 0x95)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

MARGIN = Inches(0.5)
TITLE_TOP = Inches(0.3)
TITLE_H = Inches(0.85)
CONTENT_TOP = Inches(1.25)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title(slide, text, size=32):
    box = slide.shapes.add_textbox(MARGIN, TITLE_TOP, SLIDE_W - 2 * MARGIN, TITLE_H)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = NAVY
    # thin accent rule under the title
    rule = slide.shapes.add_shape(1, MARGIN, TITLE_TOP + TITLE_H - Inches(0.05), SLIDE_W - 2 * MARGIN, Pt(2.5))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False
    return box


def add_bullets(slide, left, top, width, height, bullets, font_size=18, sub_size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ("•  " if level == 0 else "‒  ") + text
        p.level = 0
        p.font.size = Pt(font_size if level == 0 else sub_size)
        p.font.color.rgb = DARK_TEXT if level == 0 else GRAY_TEXT
        p.space_after = Pt(8 if level == 0 else 4)
    return box


def add_picture_fit(slide, path, left, top, max_w, max_h, caption=None):
    im = Image.open(path)
    iw, ih = im.size
    aspect = iw / ih
    box_aspect = max_w / max_h
    if aspect > box_aspect:
        w = max_w
        h = Emu(int(max_w / aspect))
    else:
        h = max_h
        w = Emu(int(max_h * aspect))
    x = Emu(int(left + (max_w - w) / 2))
    y = Emu(int(top + (max_h - h) / 2))
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if caption:
        cap_top = Emu(int(top + max_h + Inches(0.05)))
        box = slide.shapes.add_textbox(left, cap_top, max_w, Inches(0.3))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = GRAY_TEXT
        p.alignment = PP_ALIGN.CENTER
    return x, y, w, h


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def build():
    prs = new_presentation()

    # ---------------- Slide 1: Cover ----------------
    s = blank_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False

    title_box = s.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.33), Inches(1.3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Autonomous Driving Project"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub_box = s.shapes.add_textbox(Inches(1.0), Inches(3.7), Inches(11.33), Inches(0.6))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = 'TUM "Introduction to ROS" 2026 — Autonomous Driving Group Project'
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xC9, 0xD6, 0xE8)
    p.alignment = PP_ALIGN.CENTER

    team_box = s.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(11.33), Inches(0.5))
    tf = team_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Author 1  ·  Author 2  ·  Author 3"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0x9B, 0xB0, 0xCC)
    p.alignment = PP_ALIGN.CENTER

    set_notes(s,
        "Good afternoon. This is our autonomous driving project for the Introduction "
        "to ROS course. I will walk through our system architecture, then perception, "
        "planning, and control, and finish with our test results."
    )

    # ---------------- Slide 2: System Overview ----------------
    s = blank_slide(prs)
    add_title(s, "System Overview")
    bullets = [
        "Fixed ~785 m urban loop; waypoints published once at startup",
        "5 custom ROS 2 packages (interfaces, perception, planning, control, bringup)",
        ("plus 2 course-provided packages (simulation bridge, reference controller)", 1),
        "Perception & planning run independently, in parallel",
        "Control is the only node combining their outputs — sole path to /car_command",
        "Benchmark run: 272 s, zero collisions, zero stalls, 1 detour, 2 correct light stops",
        ("Average speed 2.7 m/s, peak 7.5 m/s on the longest straights", 1),
    ]
    add_bullets(s, MARGIN, CONTENT_TOP, Inches(6.6), Inches(5.7), bullets, font_size=16, sub_size=14)

    img_left = Inches(7.35)
    img_col_w = Inches(5.5)
    add_picture_fit(s, FIGS / "architecture_illustrated.png", img_left, CONTENT_TOP,
                     img_col_w, Inches(3.35), "Fig. 1 — Node/data-flow architecture")
    add_picture_fit(s, FIGS / "route_on_map.png", img_left, Inches(4.85),
                     img_col_w, Inches(2.1), "Fig. 5 — Route with 10 predefined goal poses")

    set_notes(s,
        "Our car drives a fixed, roughly 785-meter urban loop, using waypoints "
        "published once at startup. We wrote five ROS 2 packages -- interfaces, "
        "perception, planning, control, and bringup -- on top of two packages "
        "provided by the course. As Figure 1 shows, perception and planning run "
        "independently and in parallel, each consuming the simulator's sensor and "
        "pose streams. Control is the only node that combines their outputs, and "
        "it's the sole path to the car's command topic -- no other node talks to "
        "the simulator directly. On the results side, a full benchmark run "
        "completed in 272 seconds with zero collisions, zero stalls, one obstacle "
        "detour, and two red-light stops handled correctly. Average speed was 2.7 "
        "meters per second, peaking at 7.5 on the two longest straights. Figure 5 "
        "shows the planned route on the course map, along with the ten predefined "
        "goal poses our planner selects between."
    )

    # ---------------- Slide 3: Perception ----------------
    s = blank_slide(prs)
    add_title(s, "Perception")
    bullets = [
        "obstacle_guard_node: depth point cloud → world frame",
        ("publishes 4 corridor distances: main, tight, overtake-left, overtake-right", 1),
        ("keeps points 0.35–3.0 m high, filtering road-surface noise", 1),
        ("OctoMap cross-check (0.08–0.35 m band) catches curb-height misses", 1),
        "traffic_light_node: HSV blob detection, shape-filtered",
        ("detects red & green only; amber can't be separated from casing", 1),
        ("\"no fresh green\" is treated as \"keep waiting\"", 1),
        "No node ever uses the semantic camera — full bonus criterion claimed",
        "Limitation: ~0.7 Hz update rate, capped by simulator's single-core loop",
    ]
    add_bullets(s, MARGIN, CONTENT_TOP, SLIDE_W - 2 * MARGIN, Inches(5.9), bullets, font_size=19, sub_size=16)

    set_notes(s,
        "Perception has two nodes. obstacle_guard_node takes the depth camera's "
        "point cloud, transforms it into the world frame, and measures distance to "
        "obstacles along four corridors: main, tight, and two overtake corridors. "
        "It keeps points between 0.35 and 3 meters high, filtering out road-surface "
        "noise. It's cross-checked against an OctoMap occupancy layer covering the "
        "8-to-35-centimeter band, which catches low, curb-height obstacles the live "
        "scan structurally can't see. traffic_light_node uses HSV color-blob "
        "detection on the facing signal, filtered by shape to reject false "
        "positives like red signage. It only detects red and green -- amber can't "
        "be reliably separated from the amber-painted signal casing -- so the "
        "controller simply treats 'no fresh green yet' as 'keep waiting.' "
        "Importantly, no node in our pipeline ever uses the semantic camera, which "
        "claims the assignment's full bonus. One limitation: point-cloud updates "
        "run at only about 0.7 hertz, bounded by the simulator's single-core render "
        "loop, which caps how fast we can safely cruise."
    )

    # ---------------- Slide 4: Planning ----------------
    s = blank_slide(prs)
    add_title(s, "Planning")
    bullets = [
        "route_planner_node — two responsibilities",
        "1. Base path (once, at startup):",
        ("right-hand lane offset, so the car keeps its own side of the road", 1),
        ("iterative smoothing + curvature/acceleration-limited speed profile", 1),
        ("covers both the path-planner and trajectory-planner roles", 1),
        "2. Online updates:",
        ("detour replanning around confirmed static obstacles", 1),
        ("located via binary search on arc length, not full rescans", 1),
        ("smooth blend in/out; reverts once obstacle is clear", 1),
        ("continuous nearest-still-ahead goal selection over 10 fixed poses", 1),
    ]
    add_bullets(s, MARGIN, CONTENT_TOP, SLIDE_W - 2 * MARGIN, Inches(5.9), bullets, font_size=19, sub_size=16)

    set_notes(s,
        "Planning is a single node, route_planner_node, with two jobs. First, once "
        "at startup, it builds a base path: a right-hand lane offset so the car "
        "keeps to its own side of the road, iterative smoothing, and a curvature- "
        "and acceleration-limited speed profile -- slower through tight turns, "
        "respecting braking and acceleration limits between points. This covers "
        "both the geometric path-planning role and the kinematic trajectory-"
        "planning role. Second, online, it handles two things. When control "
        "confirms a static obstacle, the planner shifts a short stretch of the "
        "path sideways, blends smoothly back in, and reverts once clear -- "
        "locating that stretch with a binary search over arc length rather than "
        "rescanning the whole route. It also continuously selects the nearest "
        "still-ahead goal from the task's ten predefined poses, using a "
        "heading-locked search to handle the fact that our loop crosses near "
        "itself."
    )

    # ---------------- Slide 5: Control ----------------
    s = blank_slide(prs)
    add_title(s, "Control")
    bullets = [
        "pure_pursuit_node: pure-pursuit steering + proportional/feed-forward throttle",
        "/control/enable service — required, enabled by default at launch",
        "ACC-style gap control:",
        ("13 m standoff comfort-stop curve in the tight corridor", 1),
        ("14–38 m wide-corridor speed ramp; also absorbs Event I (NPC merge)", 1),
        "Traffic-light stop/go state machine",
        "Obstacle detour state machine (verify-before-commit)",
        "Emergency stop: brakes if required deceleration > 4.0 m/s² (Event II)",
    ]
    add_bullets(s, MARGIN, CONTENT_TOP, SLIDE_W - 2 * MARGIN, Inches(5.9), bullets, font_size=19, sub_size=16)

    set_notes(s,
        "Control lives in pure_pursuit_node. Steering uses a standard pure-pursuit "
        "law with a speed-dependent look-ahead distance; throttle combines "
        "proportional correction with a feed-forward table, since proportional "
        "control alone settles well below the commanded speed. The node exposes "
        "the required control-enable service, and it's enabled by default at "
        "launch. For safety, an ACC-style gap controller enforces a comfort-stop "
        "curve with a 13-meter standoff in the tight corridor, plus a "
        "wide-corridor speed ramp between 14 and 38 meters -- this same following "
        "logic absorbs the assignment's Event I, a vehicle merging in ahead. A "
        "traffic-light state machine stops and releases correctly at each signal. "
        "A detour state machine verifies a path is actually clear before "
        "committing to it. And independent of all that, every control tick checks "
        "whether emergency braking is needed -- decelerating at over 4 meters per "
        "second squared -- which is what handles Event II, a vehicle that crosses "
        "and brakes hard ahead of us."
    )

    # ---------------- Slide 6: Results ----------------
    s = blank_slide(prs)
    add_title(s, "Results")
    bullets = [
        "Fastest: two longest clear straights, peak 7.5 m/s",
        "Slowest (near-zero dips): detour spot, TL2/TL3 stops, 2 ACC slowdowns",
        ("ACC dips = Event I (merge) and Event II (crossing + hard brake, s≈531)", 1),
        "Design choice that mattered: wide-corridor cap 2.5 → 2.5–6.5 m/s ramp",
        ("roadside furniture occupies that corridor 60–90% of route time", 1),
        ("fixed route-wide starvation, not just one segment", 1),
        "Extra compute is localized to the detour segment (replan + re-profile)",
        "System-wide limiter: perception's ~0.7 Hz update rate bounds top speed",
    ]
    add_bullets(s, MARGIN, CONTENT_TOP, Inches(6.6), Inches(5.7), bullets, font_size=16, sub_size=14)

    img_left = Inches(7.35)
    img_col_w = Inches(5.5)
    add_picture_fit(s, FIGS / "route_speed_map.png", img_left, CONTENT_TOP,
                     img_col_w, Inches(3.1), "Fig. 3 — Driven path colored by speed")
    add_picture_fit(s, FIGS / "speed_profile.png", img_left, Inches(4.6),
                     img_col_w, Inches(2.3), "Fig. 4 — Speed vs. distance travelled")

    set_notes(s,
        "Figure 3 shows our driven path colored by speed, and Figure 4 shows that "
        "same speed against distance travelled. The fastest segments are the two "
        "longest clear straights, where we peak at 7.5 meters per second. The "
        "near-zero dips mark where we're slower: the detour location, the two "
        "traffic lights that were red on this run, and two ACC-managed slowdowns "
        "for the assignment's Event one and Event two encounters. One design "
        "choice mattered a lot here: we originally capped wide-corridor speed at a "
        "flat 2.5 meters per second, but roadside furniture sits in that corridor "
        "sixty to ninety percent of route time, so it was starving the whole run. "
        "Ramping that cap up to 6.5 meters per second fixed it route-wide, not "
        "just locally. The other limiter is perception's update rate -- capped "
        "around 0.7 hertz by the simulator's single-core loop -- which is what "
        "ultimately bounds how much faster we could safely go."
    )

    out_path = DOCS / "Report.pptx"
    prs.save(out_path)
    print(f"Saved {out_path}")


if __name__ == "__main__":
    build()
